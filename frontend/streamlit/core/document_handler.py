"""
Document Processing and Management for RAG LlamaStack Application
Handles file uploads, processing, and document library management
"""

import streamlit as st
import os
import time
import tempfile
import signal
import gc
from typing import List, Any, Optional
from pathlib import Path
import pandas as pd
from .utils import format_file_size, mark_upload_start, mark_upload_success, mark_upload_failed
from .web_content_processor import WebContentProcessor

# Set environment variable to avoid tokenizers parallelism warning
os.environ['TOKENIZERS_PARALLELISM'] = 'false'

# Import other modules
try:
    import docling
    DOCLING_AVAILABLE = True
except ImportError:
    DOCLING_AVAILABLE = False


# Global embedding model cache
_embedding_model = None

def read_file_content(uploaded_file) -> str:
    """Extract text content from uploaded files with performance optimization and robust error handling"""
    try:
        # Read file content efficiently
        file_content = uploaded_file.read()
        file_size_mb = len(file_content) / (1024 * 1024)
        
        print(f"📄 Processing {uploaded_file.name} ({file_size_mb:.1f}MB)")
        
        # Optimize extraction based on file size
        if file_size_mb > 5:  # Large files need optimization
            print(f"🚀 Optimizing extraction for large file ({file_size_mb:.1f}MB)")
        
        # Try docling for all supported file types with robust error handling
        content = try_docling_extraction_enhanced(uploaded_file, file_content, file_size_mb)
        if content:
            return content
        
        # Fallback to specialized extraction for each file type
        content = try_specialized_extraction(uploaded_file, file_content, file_size_mb)
        if content:
            return content
        
        # Final fallback to basic text extraction
        content = try_basic_extraction(uploaded_file, file_content)
        if content:
            return content
        
        return ""
        
    except Exception as e:
        print(f"❌ Failed to read {uploaded_file.name}: {e}")
        return ""

def try_docling_extraction_enhanced(uploaded_file, file_content: bytes, file_size_mb: float) -> str:
    """Try docling extraction for all supported file types with robust error handling"""
    try:
        from docling.document_converter import DocumentConverter
        
        # Create converter with optimized settings for large files
        converter = DocumentConverter()
        
        # Determine file extension and MIME type
        file_extension = get_file_extension(uploaded_file.name)
        mime_type = uploaded_file.type
        
        # Map file types to appropriate suffixes
        suffix_map = {
            'pdf': '.pdf',
            'docx': '.docx', 
            'pptx': '.pptx',
            'txt': '.txt',
            'md': '.md'
        }
        
        suffix = suffix_map.get(file_extension.lower(), '.txt')
        
        # Use a more secure temporary file creation
        import tempfile
        import os
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix, mode='wb') as tmp_file:
            try:
                tmp_file.write(file_content)
                tmp_file.flush()
                tmp_file_path = tmp_file.name
            except Exception as write_error:
                print(f"⚠️ Failed to write temporary file: {write_error}")
                return None
        
        try:
            # Add timeout and memory protection
            import signal
            
            def timeout_handler(signum, frame):
                raise TimeoutError(f"{file_extension.upper()} processing timed out")
            
            # Set a timeout based on file size (larger files get more time)
            timeout_seconds = min(60, max(15, int(file_size_mb * 10)))  # 15-60 seconds
            signal.signal(signal.SIGALRM, timeout_handler)
            signal.alarm(timeout_seconds)
            
            try:
                print(f"🔧 Converting {file_extension.upper()} with docling (timeout: {timeout_seconds}s)...")
                result = converter.convert(tmp_file_path)
                content = result.document.export_to_markdown()
                signal.alarm(0)  # Cancel the alarm
                
                # Clean up temp file
                try:
                    os.unlink(tmp_file_path)
                except:
                    pass
                
                if content and content.strip():
                    print(f"✅ Extracted {len(content)} chars from {uploaded_file.name} using docling")
                    
                    # For very large content, do preliminary optimization
                    if len(content) > 500000:  # 500KB of text
                        content = optimize_extracted_content(content)
                        print(f"🚀 Optimized content to {len(content)} chars")
                        
                    return content
                else:
                    print("⚠️ Docling returned empty content")
                    return None
                    
            except TimeoutError:
                print(f"⚠️ {file_extension.upper()} processing timed out, falling back to other methods")
                signal.alarm(0)  # Cancel the alarm
                return None
            except Exception as conversion_error:
                print(f"⚠️ Docling conversion failed for {file_extension.upper()}: {conversion_error}")
                signal.alarm(0)  # Cancel the alarm
                return None
                
        except Exception as e:
            print(f"⚠️ Docling processing error for {file_extension.upper()}: {e}")
            return None
        finally:
            # Ensure temp file is cleaned up
            try:
                if os.path.exists(tmp_file_path):
                    os.unlink(tmp_file_path)
            except:
                pass
                
    except ImportError:
        print("⚠️ Docling not available, falling back to specialized extraction")
        return None
    except Exception as e:
        print(f"⚠️ Docling extraction failed for {uploaded_file.name}: {e}")
        return None

def try_specialized_extraction(uploaded_file, file_content: bytes, file_size_mb: float) -> str:
    """Try specialized extraction methods for each file type"""
    file_extension = get_file_extension(uploaded_file.name)
    
    try:
        if file_extension.lower() == 'pdf':
            return try_multiple_pdf_extractors(file_content, uploaded_file.name)
        elif file_extension.lower() == 'docx':
            return try_docx_extraction(file_content, uploaded_file.name)
        elif file_extension.lower() == 'pptx':
            return try_pptx_extraction(file_content, uploaded_file.name)
        elif file_extension.lower() in ['txt', 'md']:
            return try_text_extraction(file_content, uploaded_file.name)
        else:
            print(f"⚠️ No specialized extractor for {file_extension}")
            return None
            
    except Exception as e:
        print(f"⚠️ Specialized extraction failed for {file_extension}: {e}")
        return None

def try_docx_extraction(file_content: bytes, filename: str) -> str:
    """Try DOCX extraction with multiple methods"""
    print(f"🔧 Trying DOCX extraction for {filename}")
    
    # Method 1: python-docx
    try:
        from docx import Document
        import io
        
        doc = Document(io.BytesIO(file_content))
        content = ""
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content += paragraph.text + "\n"
        
        if content.strip():
            print(f"✅ python-docx extracted {len(content)} chars from {filename}")
            return content
        else:
            print("⚠️ python-docx returned empty content")
            
    except ImportError:
        print("⚠️ python-docx not available")
    except Exception as e:
        print(f"⚠️ python-docx extraction failed: {e}")
    
    # Method 2: Try with docling fallback (already handled above)
    print("⚠️ DOCX extraction failed, using fallback")
    return None

def try_pptx_extraction(file_content: bytes, filename: str) -> str:
    """Try PPTX extraction with multiple methods"""
    print(f"🔧 Trying PPTX extraction for {filename}")
    
    # Method 1: python-pptx
    try:
        from pptx import Presentation
        import io
        
        prs = Presentation(io.BytesIO(file_content))
        content = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content += shape.text + "\n"
        
        if content.strip():
            print(f"✅ python-pptx extracted {len(content)} chars from {filename}")
            return content
        else:
            print("⚠️ python-pptx returned empty content")
            
    except ImportError:
        print("⚠️ python-pptx not available")
    except Exception as e:
        print(f"⚠️ python-pptx extraction failed: {e}")
    
    # Method 2: Try with docling fallback (already handled above)
    print("⚠️ PPTX extraction failed, using fallback")
    return None

def try_text_extraction(file_content: bytes, filename: str) -> str:
    """Try text extraction for TXT and MD files"""
    print(f"🔧 Trying text extraction for {filename}")
    
    try:
        # Try UTF-8 first
        content = file_content.decode('utf-8')
        if content.strip():
            print(f"✅ Text extraction successful: {len(content)} chars from {filename}")
            return content
    except UnicodeDecodeError:
        try:
            # Try with error handling
            content = file_content.decode('utf-8', errors='ignore')
            if content.strip():
                print(f"✅ Text extraction with error handling: {len(content)} chars from {filename}")
                return content
        except Exception as e:
            print(f"⚠️ Text extraction failed: {e}")
    
    print("⚠️ Text extraction failed")
    return None

def get_file_extension(filename: str) -> str:
    """Get file extension from filename"""
    return filename.split('.')[-1] if '.' in filename else ''

def try_basic_extraction(uploaded_file, file_content: bytes) -> str:
    """Try basic text extraction with multiple fallbacks"""
    try:
        if uploaded_file.type == "text/plain":
            content = file_content.decode('utf-8')
        elif uploaded_file.type == "application/pdf":
            # Try multiple PDF extraction methods
            content = try_multiple_pdf_extractors(file_content, uploaded_file.name)
        elif uploaded_file.name.endswith('.md'):
            content = file_content.decode('utf-8')
        else:
            # Try to decode as text
            try:
                content = file_content.decode('utf-8')
            except:
                content = file_content.decode('utf-8', errors='ignore')
        
        if content and content.strip():
            print(f"✅ Extracted {len(content)} chars using basic fallback method")
            
            # Optimize if content is very large
            if len(content) > 500000:
                content = optimize_extracted_content(content)
                print(f"🚀 Optimized content to {len(content)} chars")
            
            return content
        else:
            print("⚠️ Basic extraction returned empty content")
            return None
            
    except Exception as e:
        print(f"❌ Basic extraction failed: {e}")
        return None

def try_multiple_pdf_extractors(file_content: bytes, filename: str) -> str:
    """Try multiple PDF extraction methods with fallbacks"""
    
    # Method 1: Try PyPDF2
    try:
        import PyPDF2
        import io
        
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        content = ""
        
        for page_num, page in enumerate(pdf_reader.pages):
            try:
                page_text = page.extract_text()
                if page_text:
                    content += page_text + "\n"
            except Exception as page_error:
                print(f"⚠️ Failed to extract page {page_num}: {page_error}")
                continue
        
        if content.strip():
            print(f"✅ PyPDF2 extracted {len(content)} chars from {filename}")
            return content
        else:
            print("⚠️ PyPDF2 returned empty content")
            
    except ImportError:
        print("⚠️ PyPDF2 not available")
    except Exception as e:
        print(f"⚠️ PyPDF2 extraction failed: {e}")
    
    # Method 2: Try pdfplumber
    try:
        import pdfplumber
        
        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            content = ""
            for page in pdf.pages:
                try:
                    page_text = page.extract_text()
                    if page_text:
                        content += page_text + "\n"
                except Exception as page_error:
                    print(f"⚠️ Failed to extract page with pdfplumber: {page_error}")
                    continue
            
            if content.strip():
                print(f"✅ pdfplumber extracted {len(content)} chars from {filename}")
                return content
            else:
                print("⚠️ pdfplumber returned empty content")
                
    except ImportError:
        print("⚠️ pdfplumber not available")
    except Exception as e:
        print(f"⚠️ pdfplumber extraction failed: {e}")
    
    # Method 3: Simple fallback
    print("⚠️ All PDF extractors failed, using simple fallback")
    return f"[PDF content from {filename} - {len(file_content)} bytes] - PDF extraction failed. Please ensure the PDF is not corrupted or password-protected."


def optimize_extracted_content(content: str) -> str:
    """Optimize extracted content for large files by removing low-value text"""
    if len(content) <= 100000:  # Only optimize very large content
        return content
    
    lines = content.split('\n')
    optimized_lines = []
    
    for line in lines:
        line = line.strip()
        
        # Skip very short lines (likely formatting artifacts)
        if len(line) < 5:
            continue
            
        # Skip lines that are mostly numbers/symbols (likely tables/formatting)
        if len(line) > 10 and sum(c.isalpha() for c in line) / len(line) < 0.5:
            continue
            
        # Skip repetitive lines (headers/footers)
        if line.lower().startswith(('page ', 'chapter ', 'section ')):
            continue
            
        optimized_lines.append(line)
    
    optimized_content = '\n'.join(optimized_lines)
    
    # If still too large, take first 80% (usually most important content)
    if len(optimized_content) > 300000:
        optimized_content = optimized_content[:300000] + "\n\n[Content truncated for performance]"
    
    return optimized_content


def extract_pdf_text_simple(file_content: bytes) -> str:
    """Simple PDF text extraction fallback"""
    try:
        import io
        # This is a very basic fallback - in practice you'd use PyPDF2 or similar
        # For now, just return a placeholder
        return f"[PDF content extracted - {len(file_content)} bytes] - Install docling for better PDF processing"
    except Exception:
        return ""


def create_enhanced_chunks(content: str, chunk_size: int = None, overlap: int = None) -> List[str]:
    """Create optimized text chunks with intelligent splitting"""
    # Use optimized config defaults
    from .config import CHARS_PER_CHUNK, CHUNK_OVERLAP
    
    if chunk_size is None:
        chunk_size = CHARS_PER_CHUNK  # Now 3000 for better focus
    if overlap is None:
        overlap = CHUNK_OVERLAP  # Now 600
    
    if not content or len(content.strip()) == 0:
        return [""]
    
    # Clean up the content
    content = content.strip()
    
    # If content is smaller than chunk size, return as single chunk
    if len(content) <= chunk_size:
        return [content]
    
    chunks = []
    start = 0
    
    while start < len(content):
        # Calculate end position
        end = start + chunk_size
        
        # If this is the last chunk, take everything remaining
        if end >= len(content):
            chunks.append(content[start:])
            break
        
        # Try to find a good breaking point (sentence, paragraph, or word boundary)
        chunk_text = content[start:end]
        
        # Look for sentence endings within the last 200 characters
        for boundary in ['. ', '.\n', '!\n', '?\n', '. \n']:
            boundary_pos = chunk_text.rfind(boundary)
            if boundary_pos > chunk_size - 200:  # Within last 200 chars
                end = start + boundary_pos + len(boundary)
                break
        else:
            # Look for paragraph breaks
            boundary_pos = chunk_text.rfind('\n\n')
            if boundary_pos > chunk_size - 300:  # Within last 300 chars
                end = start + boundary_pos + 2
            else:
                # Look for any newline
                boundary_pos = chunk_text.rfind('\n')
                if boundary_pos > chunk_size - 100:  # Within last 100 chars
                    end = start + boundary_pos + 1
                else:
                    # Look for word boundary
                    boundary_pos = chunk_text.rfind(' ')
                    if boundary_pos > chunk_size - 50:  # Within last 50 chars
                        end = start + boundary_pos + 1
        
        # Add the chunk
        chunk = content[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        # Move start position with overlap
        start = max(start + 1, end - overlap)  # Ensure progress
    
    print(f"✂️ Created {len(chunks)} chunks (avg {sum(len(c) for c in chunks) // len(chunks)} chars each)")
    return chunks


def get_local_embedding_model():
    """Get or initialize the local sentence-transformers model"""
    global _embedding_model
    
    if _embedding_model is None:
        try:
            from sentence_transformers import SentenceTransformer
            # Use a fast, lightweight model that works well for RAG
            _embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
            st.session_state.embedding_model_loaded = True
            st.session_state.embedding_model_name = 'all-MiniLM-L6-v2'
        except ImportError:
            st.error("❌ sentence-transformers not installed. Run: pip install sentence-transformers")
            return None
        except Exception as e:
            st.warning(f"⚠️ Could not load embedding model: {e}")
            return None
    
    return _embedding_model


def generate_embeddings_batch_optimized(chunks: List[str], progress_bar, status_text, start_progress: float = 0.0, file_id: str = None) -> tuple[List[List[float]], int]:
    """Generate embeddings with progress tracking and error handling"""
    if not chunks:
        return [], 0
    
    model = get_local_embedding_model()
    embeddings = []
    embedding_errors = 0
    
    try:
        # Generate real embeddings in batches for efficiency
        batch_size = 32
        total_batches = (len(chunks) + batch_size - 1) // batch_size
        
        for i in range(0, len(chunks), batch_size):
            # Check for interruption
            if file_id and file_id not in st.session_state.currently_uploading:
                raise InterruptedError(f"Upload interrupted for {file_id}")
            
            batch = chunks[i:i + batch_size]
            batch_num = i // batch_size + 1
            
            # Update progress
            progress = start_progress + (batch_num / total_batches) * 0.4  # 40% of remaining progress
            progress_bar.progress(progress)
            status_text.text(f"🧠 Generating embeddings... ({batch_num}/{total_batches})")
            
            try:
                if model is not None:
                    # Use real model
                    batch_embeddings = model.encode(batch, convert_to_tensor=False, show_progress_bar=False)
                    embeddings.extend(batch_embeddings.tolist())
                else:
                    # Fallback to dummy embeddings
                    for j, text in enumerate(batch):
                        import random
                        dummy_embedding = [random.uniform(-0.1, 0.1) for _ in range(384)]
                        # Add variation based on text content
                        text_hash = hash(text) % 1000
                        for k in range(min(10, len(dummy_embedding))):
                            dummy_embedding[k] += (text_hash / 10000.0) + ((i + j) / 1000.0)
                        embeddings.append(dummy_embedding)
                        
            except Exception as e:
                print(f"⚠️ Error generating embeddings for batch {batch_num}: {e}")
                embedding_errors += len(batch)
                # Add dummy embeddings for failed batch
                for _ in batch:
                    import random
                    dummy_embedding = [random.uniform(-0.1, 0.1) for _ in range(384)]
                    embeddings.append(dummy_embedding)
        
        return embeddings, embedding_errors
        
    except InterruptedError:
        raise
    except Exception as e:
        st.warning(f"⚠️ Embedding generation failed: {e}")
        # Fallback to all dummy embeddings
        import random
        embeddings = []
        for i, text in enumerate(chunks):
            dummy_embedding = [random.uniform(-0.1, 0.1) for _ in range(384)]
            text_hash = hash(text) % 1000
            for j in range(min(10, len(dummy_embedding))):
                dummy_embedding[j] += (text_hash / 10000.0) + (i / 1000.0)
            embeddings.append(dummy_embedding)
        return embeddings, len(chunks)  # All failed


def generate_real_embeddings(texts: List[str]) -> List[List[float]]:
    """Generate real embeddings using local sentence-transformers"""
    model = get_local_embedding_model()
    
    if model is None:
        # Fallback to varied dummy embeddings
        import random
        embeddings = []
        for i, text in enumerate(texts):
            dummy_embedding = [random.uniform(-0.1, 0.1) for _ in range(384)]
            # Add variation based on text content
            text_hash = hash(text) % 1000
            for j in range(min(10, len(dummy_embedding))):
                dummy_embedding[j] += (text_hash / 10000.0) + (i / 1000.0)
            embeddings.append(dummy_embedding)
        return embeddings
    
    try:
        # Generate real embeddings in batches for efficiency
        batch_size = 32
        all_embeddings = []
        
        for i in range(0, len(texts), batch_size):
            batch = texts[i:i + batch_size]
            batch_embeddings = model.encode(batch, convert_to_tensor=False, show_progress_bar=False)
            all_embeddings.extend(batch_embeddings.tolist())
        
        return all_embeddings
    except Exception as e:
        st.warning(f"⚠️ Embedding generation failed: {e}")
        # Fallback to dummy embeddings
        import random
        embeddings = []
        for i, text in enumerate(texts):
            dummy_embedding = [random.uniform(-0.1, 0.1) for _ in range(384)]
            text_hash = hash(text) % 1000
            for j in range(min(10, len(dummy_embedding))):
                dummy_embedding[j] += (text_hash / 10000.0) + (i / 1000.0)
            embeddings.append(dummy_embedding)
def backup_documents_data():
    """Backup document data to prevent loss during reruns"""
    if 'uploaded_documents' in st.session_state and st.session_state.uploaded_documents:
        # Store in a more persistent way
        st.session_state['_backup_uploaded_documents'] = st.session_state.uploaded_documents.copy()
        st.session_state['_backup_documents'] = st.session_state.documents.copy() if 'documents' in st.session_state else []
        print(f"💾 Backed up {len(st.session_state.uploaded_documents)} documents")

def restore_documents_data():
    """Restore document data if it was lost during rerun"""
    restored = False
    
    # Restore uploaded_documents if lost
    if ('uploaded_documents' not in st.session_state or not st.session_state.uploaded_documents) and '_backup_uploaded_documents' in st.session_state:
        st.session_state.uploaded_documents = st.session_state['_backup_uploaded_documents'].copy()
        print(f"🔄 Restored {len(st.session_state.uploaded_documents)} uploaded documents from backup")
        restored = True
    
    # Restore documents if lost  
    if ('documents' not in st.session_state or not st.session_state.documents) and '_backup_documents' in st.session_state:
        st.session_state.documents = st.session_state['_backup_documents'].copy()
        print(f"🔄 Restored {len(st.session_state.documents)} documents from backup")
        restored = True
    
    return restored


def initialize_document_storage() -> None:
    """Initialize document storage in session state with backup restoration"""
    # Try to restore from backup first
    restore_documents_data()
    
    # Initialize if still empty
    if "uploaded_documents" not in st.session_state:
        st.session_state.uploaded_documents = []
        print("📋 Initialized empty uploaded_documents")
    if 'documents' not in st.session_state:
        st.session_state.documents = []
        print("📄 Initialized empty documents")
    
    # Create initial backup
    backup_documents_data()


def render_file_uploader() -> List:
    """Render file upload interface with original layout"""
    # Show current upload status
    if hasattr(st.session_state, 'currently_uploading') and st.session_state.currently_uploading:
        uploading_files = [file_id.split('_')[0] for file_id in st.session_state.currently_uploading]
        st.warning(f"⏳ Currently uploading: {', '.join(uploading_files)}")
        st.warning("⚠️ Don't change models during upload!")
    
    if hasattr(st.session_state, 'failed_uploads') and st.session_state.failed_uploads:
        failed_count = len(st.session_state.failed_uploads)
        st.error(f"🔄 {failed_count} file(s) need retry due to interruption")
        if st.button("🗑️ Clear failed uploads", type="secondary"):
            st.session_state.failed_uploads.clear()
            st.rerun()
    
    # Input method selector
    if 'input_method' not in st.session_state:
        st.session_state.input_method = "📄 Upload Files"
    
    input_method = st.radio(
        "Choose input method:",
        ["📄 Upload Files", "🌐 Process Web URLs", "📂 Upload URLs File"],
        key="input_method_radio",
        horizontal=True
    )
    
    uploaded_files = []
    
    if input_method == "📄 Upload Files":
        uploaded_files = st.file_uploader(
            "Choose files to upload",
            type=['pdf', 'txt', 'md', 'docx', 'pptx'],
            accept_multiple_files=True,
            help="Drag and drop files here (Max 50MB per file)"
        )
        
        # Upload tips in collapsible
        with st.expander("💡 Upload Tips", expanded=False):
            st.markdown("""
            **Best Practices:**
            - ✅ Upload one file at a time for large files (>10MB)
            - ✅ Don't switch models during upload
            - ✅ Wait for "Processing complete!" before making changes
            - ✅ Failed uploads can be retried automatically
            
            **Supported Formats:**
            - 📄 PDF documents
            - 📝 Text files (.txt, .md)
            - 📘 Word documents (.docx)
            - 📊 PowerPoint (.pptx)
            """)
    
    elif input_method == "🌐 Process Web URLs":
        # Call the bulk URL processing interface
        render_url_paste_input()

    elif input_method == "📂 Upload URLs File":
        render_url_file_upload_input()
    
    return uploaded_files if uploaded_files else []


def process_web_url(url: str):
    """Process a web URL and add to document storage"""
    if not url.strip():
        st.error("Please enter a valid URL")
        return
    
    try:
        # Initialize web content processor if not exists
        if 'web_content_processor' not in st.session_state:
            st.session_state.web_content_processor = WebContentProcessor()
        
        # Get web content processor from session state
        processor = st.session_state.web_content_processor
        
        # Extract content from URL
        result = processor.process_url(url)
        
        if result and result.get('success'):
            # Create a document-like structure for web content
            content = result.get('content', '')
            title = result.get('title', f"Web Content from {url}")
            
            # Get domain for display
            from urllib.parse import urlparse
            domain = urlparse(url).netloc
            
            web_document = {
                'name': f"{title}",
                'content': content,
                'file_type': 'WEB',
                'file_size_mb': len(content) / (1024 * 1024),
                'chunk_count': 0,  # Will be calculated during embedding
                'character_count': len(content),
                'upload_time': time.strftime("%Y-%m-%d %H:%M:%S"),
                'source_url': url,
                'domain': domain,
                'extraction_method': result.get('method', 'unknown'),
                'metadata': result.get('metadata', {}),
                'processed_at': result.get('processed_at', time.time())
            }
            
            # Process content for embeddings using existing embedding system
            with st.spinner("🧮 Creating embeddings for web content..."):
                # Create chunks using existing chunking logic
                chunks = create_chunks_from_content(content)
                web_document['chunk_count'] = len(chunks)
                
                # Create embeddings using existing embedding system
                embeddings_data = create_embeddings_from_chunks(chunks)
                
                if embeddings_data:
                    web_document['chunks'] = chunks
                    web_document['embeddings'] = embeddings_data
                    
                    # Add to uploaded documents using existing storage
                    if 'uploaded_documents' not in st.session_state:
                        st.session_state.uploaded_documents = []
                    
                    # Check for duplicates before adding
                    existing_web_docs = [doc for doc in st.session_state.uploaded_documents 
                                       if doc.get('source_url') == url]
                    
                    if existing_web_docs:
                        st.warning(f"⚠️ Content from '{domain}' already exists in your documents")
                        return
                    
                    st.session_state.uploaded_documents.append(web_document)
                    
                    # Show success message
                    st.success(f"✅ Successfully processed web content from {domain}")
                    
                    # Show MCP server information
                    if doc.get('extraction_method'):
                        extraction_method = doc.get('extraction_method')
                        method_map = {
                            'mcp_just-every': '📦 Just-Every MCP'
                        }
                        method_display = method_map.get(extraction_method, extraction_method)
                        st.info(f"🔧 **MCP Server Used:** {method_display}")
                    
                    # Show processing stats
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("📝 Characters", f"{web_document['character_count']:,}")
                    with col2:
                        st.metric("🧩 Chunks", web_document['chunk_count'])
                    with col3:
                        st.metric("📊 Size", f"{web_document['file_size_mb']:.2f} MB")
                    
                    # Form will automatically clear the input on successful submission
                else:
                    st.error("❌ Failed to create embeddings for web content")
        else:
            st.error("❌ Failed to extract content from the URL. Please check the URL and try again.")
            
    except Exception as e:
        st.error(f"❌ Error processing URL: {str(e)}")


def get_memory_usage():
    """Get current memory usage in MB"""
    try:
        import psutil
        process = psutil.Process(os.getpid())
        return process.memory_info().rss / 1024 / 1024  # Convert to MB
    except ImportError:
        return 0

def process_bulk_web_urls(urls: List[str], progress_callback=None):
    """Process multiple web URLs in bulk and add to document storage with improved memory management"""
    if not urls:
        st.error("Please provide at least one URL")
        return
    
    # Import urlparse at the top of the function
    from urllib.parse import urlparse
    
    # Limit the number of URLs to prevent overwhelming the system
    MAX_URLS = 25  # Reduced to prevent memory issues
    if len(urls) > MAX_URLS:
        st.warning(f"⚠️ Too many URLs ({len(urls)}). Processing first {MAX_URLS} URLs only.")
        urls = urls[:MAX_URLS]
    
    # Clean and validate URLs
    valid_urls = []
    invalid_urls = []
    for url in urls:
        url = url.strip()
        if url and url.startswith(('http://', 'https://')):
            valid_urls.append(url)
        elif url:
            # Try to add https if no protocol specified
            if not url.startswith(('http://', 'https://')):
                url = f"https://{url}"
                valid_urls.append(url)
        else:
            invalid_urls.append(url)
    
    # Show URL validation results
    st.info(f"📊 URL Validation Results:")
    st.markdown(f"""
    - 📄 **Total URLs in file**: {len(urls)}
    - ✅ **Valid URLs**: {len(valid_urls)}
    - ❌ **Invalid URLs**: {len(invalid_urls)}
    """)
    
    if invalid_urls:
        with st.expander("❌ Invalid URLs (will be skipped)"):
            for invalid_url in invalid_urls:
                st.code(invalid_url)
    
    if not valid_urls:
        st.error("No valid URLs found. Please provide URLs starting with http:// or https://")
        return
    
    # Show processing start message
    st.info(f"🚀 Starting bulk processing of {len(valid_urls)} URLs...")
    
    # Monitor initial memory usage
    initial_memory = get_memory_usage()
    if initial_memory > 0:
        st.info(f"💾 Initial memory usage: {initial_memory:.1f} MB")
    
    # Initialize web content processor if not exists
    if 'web_content_processor' not in st.session_state:
        st.session_state.web_content_processor = WebContentProcessor()
    
    processor = st.session_state.web_content_processor
    
    # Initialize session state for uploaded documents
    if 'uploaded_documents' not in st.session_state:
        st.session_state.uploaded_documents = []
    
    # Track processing results
    results = {
        'successful': [],
        'failed': [],
        'duplicates': [],
        'total_processed': 0
    }
    
    # Create progress bar for bulk processing
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    try:
        for i, url in enumerate(valid_urls):
            # Update progress
            progress = (i + 1) / len(valid_urls)
            progress_bar.progress(progress)
            status_text.text(f"Processing {i + 1}/{len(valid_urls)}: {url}")
            
            try:
                # Extract content from URL using silent mode to reduce alerts
                result = processor.process_url(url, silent_mode=True)
                
                if result and result.get('success'):
                    # Create a document-like structure for web content
                    content = result.get('content', '')
                    title = result.get('title', f"Web Content from {url}")
                    
                    # Get domain for display
                    from urllib.parse import urlparse
                    domain = urlparse(url).netloc
                    
                    # Check for duplicates before processing
                    existing_web_docs = [doc for doc in st.session_state.uploaded_documents 
                                       if doc.get('source_url') == url]
                    
                    if existing_web_docs:
                        results['duplicates'].append({
                            'url': url,
                            'domain': domain,
                            'reason': 'Already exists in documents'
                        })
                        continue
                    
                    web_document = {
                        'name': f"{title}",
                        'content': content,
                        'file_type': 'WEB',
                        'file_size_mb': len(content) / (1024 * 1024),
                        'chunk_count': 0,  # Will be calculated during embedding
                        'character_count': len(content),
                        'upload_time': time.strftime("%Y-%m-%d %H:%M:%S"),
                        'source_url': url,
                        'domain': domain
                    }
                    
                    # Create chunks and embeddings for web content
                    try:
                        # Update status to show chunking
                        status_text.text(f"Processing {i + 1}/{len(valid_urls)}: {url} - Creating chunks...")
                        
                        # Create chunks using existing chunking logic
                        chunks = create_chunks_from_content(content)
                        web_document['chunk_count'] = len(chunks)
                        
                        # Update status to show embedding creation
                        status_text.text(f"Processing {i + 1}/{len(valid_urls)}: {url} - Creating embeddings...")
                        
                        # Create embeddings using existing embedding system
                        embeddings_data = create_embeddings_from_chunks(chunks)
                        
                        if embeddings_data:
                            web_document['chunks'] = chunks
                            web_document['embeddings'] = embeddings_data
                            
                            # Store embeddings in VectorIO database
                            vectorio_success = False
                            try:
                                if 'llamastack_client' in st.session_state:
                                    # Prepare metadata for VectorIO
                                    vector_metadata = []
                                    for j, chunk in enumerate(chunks):
                                        vector_metadata.append({
                                            "document_id": f"web_{domain}_{j}",
                                            "document_name": f"Web Content from {url}",
                                            "chunk_index": j,
                                            "chunk_content": chunk,
                                            "file_type": "WEB",
                                            "source_url": url,
                                            "domain": domain,
                                            "title": title,
                                            "upload_time": time.strftime("%Y-%m-%d %H:%M:%S"),
                                            "chunk_count": len(chunks),
                                            "character_count": len(content)
                                        })
                                    
                                    vectorio_success = st.session_state.llamastack_client.store_embeddings_in_vector_db(
                                        embeddings=embeddings_data,
                                        metadata=vector_metadata,
                                        vector_db_id="faiss"
                                    )
                                    if vectorio_success:
                                        print(f"✅ Successfully stored web content embeddings in VectorIO for {url}")
                                    else:
                                        print(f"⚠️ Failed to store web content in VectorIO for {url}")
                            except Exception as e:
                                print(f"⚠️ VectorIO storage failed for web content {url}: {e}")
                            
                            web_document['vectorio_stored'] = vectorio_success
                            
                            # Add to session state
                            st.session_state.uploaded_documents.append(web_document)
                            results['successful'].append({
                                'url': url,
                                'domain': domain,
                                'title': title,
                                'content_length': len(content),
                                'chunks': len(chunks),
                                'vectorio_stored': vectorio_success
                            })
                            
                            # Force garbage collection after each successful processing
                            gc.collect()
                            # Clear any large objects from memory
                            if 'web_document' in locals():
                                del web_document
                            if 'content' in locals():
                                del content
                            if 'chunks' in locals():
                                del chunks
                            if 'embeddings_data' in locals():
                                del embeddings_data
                            gc.collect()
                        else:
                            # If embedding creation failed, still store the document but mark it
                            web_document['embedding_error'] = 'Failed to create embeddings'
                            st.session_state.uploaded_documents.append(web_document)
                            results['successful'].append({
                                'url': url,
                                'domain': domain,
                                'title': title,
                                'content_length': len(content),
                                'chunks': len(chunks),
                                'note': 'Content stored but embeddings failed'
                            })
                    except Exception as embedding_error:
                        # If chunking/embedding fails, still store the document
                        web_document['embedding_error'] = str(embedding_error)
                        st.session_state.uploaded_documents.append(web_document)
                        results['successful'].append({
                            'url': url,
                            'domain': domain,
                            'title': title,
                            'content_length': len(content),
                            'chunks': 0,
                            'note': f'Content stored but processing failed: {str(embedding_error)}'
                        })
                    
                else:
                    # Handle failed processing - result is None or success is False
                    if result is None:
                        error_msg = "No response from processor - all extraction methods failed"
                    else:
                        error_msg = result.get('error', 'Unknown error during processing')
                    
                    results['failed'].append({
                        'url': url,
                        'domain': urlparse(url).netloc if url else 'unknown',
                        'error': error_msg
                    })
                    
            except Exception as e:
                # Handle individual URL processing errors
                error_msg = str(e)
                results['failed'].append({
                    'url': url,
                    'domain': urlparse(url).netloc if url else 'unknown',
                    'error': error_msg
                })
                
                # Continue with next URL instead of stopping
                continue
            
            # Add a small delay between requests to prevent overwhelming servers
            time.sleep(0.5)
            
            # Force garbage collection periodically
            if (i + 1) % 5 == 0:
                gc.collect()
                # Monitor memory usage every 5 URLs
                current_memory = get_memory_usage()
                if current_memory > 0:
                    memory_increase = current_memory - initial_memory
                    if memory_increase > 500:  # If memory increased by more than 500MB
                        st.warning(f"⚠️ High memory usage detected: {current_memory:.1f} MB (+{memory_increase:.1f} MB)")
                        # Force more aggressive cleanup
                        gc.collect()
                        gc.collect()
                        # Add a small delay to allow memory cleanup
                        time.sleep(1)
        
        # Update final progress
        progress_bar.progress(1.0)
        status_text.text("✅ Bulk processing completed!")
        
        # Show results summary
        st.success(f"✅ Bulk processing completed!")
        
        # Report final memory usage
        final_memory = get_memory_usage()
        if final_memory > 0 and initial_memory > 0:
            memory_change = final_memory - initial_memory
            st.info(f"💾 Final memory usage: {final_memory:.1f} MB (change: {memory_change:+.1f} MB)")
        
        # Calculate total chunks and content
        total_chunks = sum(item.get('chunks', 0) for item in results['successful'])
        total_content_length = sum(item.get('content_length', 0) for item in results['successful'])
        
        st.markdown(f"""
        **Processing Results:**
        - ✅ **Successful**: {len(results['successful'])} URLs
        - ❌ **Failed**: {len(results['failed'])} URLs
        - 🔄 **Duplicates**: {len(results['duplicates'])} URLs
        - 📊 **Total Processed**: {len(valid_urls)} URLs
        - 📄 **Total Chunks Created**: {total_chunks} chunks
        - 📝 **Total Content**: {total_content_length:,} characters
        """)
        
        # Show successful URLs with their processing details
        if results['successful']:
            with st.expander(f"✅ {len(results['successful'])} URLs processed successfully"):
                for item in results['successful']:
                    st.markdown(f"**{item['domain']}** - {item['title']}")
                    st.caption(f"📄 {item['chunks']} chunks | 📝 {item['content_length']:,} characters")
                    if 'note' in item:
                        st.caption(f"⚠️ {item['note']}")
                    st.markdown("---")
        
        # Show failed URLs if any
        if results['failed']:
            with st.expander("❌ Failed URLs"):
                for failed in results['failed']:
                    st.markdown(f"**{failed['url']}** - {failed['error']}")
        
        # Show duplicate URLs if any
        if results['duplicates']:
            with st.expander("�� Duplicate URLs"):
                for duplicate in results['duplicates']:
                    st.markdown(f"**{duplicate['url']}** - {duplicate['reason']}")
        
        # Final garbage collection
        gc.collect()
        
    except Exception as e:
        st.error(f"❌ Error during bulk processing: {str(e)}")
        # Clean up on error
        gc.collect()
    finally:
        # Always clean up resources
        gc.collect()


def render_url_paste_input():
    """Render the URL paste input interface"""
    st.markdown("Process web URLs to extract content. Enter one URL per line (or just one URL for single processing):")
    
    # Text area for bulk URLs
    urls_text = st.text_area(
        "Enter URLs (one per line):",
        placeholder="https://github.com/username/repo/blob/main/README.md\nhttps://example.com\nhttps://docs.example.com",
        height=150,
        help="Enter one or multiple URLs, one per line. URLs should start with http:// or https://",
        key="bulk_urls_text"
    )

    # Process manually entered URLs
    if urls_text:
        urls = [url.strip() for url in urls_text.split('\n') if url.strip()]
        
        if urls:
            st.markdown(f"**Found {len(urls)} URLs to process:**")
            for url in urls[:5]:  # Show first 5
                st.code(url)
            if len(urls) > 5:
                st.caption(f"... and {len(urls) - 5} more URLs")
            
            if st.button("🚀 Process URLs", type="primary"):
                process_bulk_web_urls(urls)
        else:
            st.warning("⚠️ No valid URLs found in the text area")
    else:
        st.info("💡 Enter URLs in the text area above to process them.")

    # Show processing tips
    with st.expander("💡 Processing Tips"):
        st.markdown("""
        **Best Practices for URL Processing:**
        
        - **Single URL**: Just enter one URL for individual processing
        - **Multiple URLs**: Enter multiple URLs, one per line
        - **URL Format**: Ensure URLs start with `http://` or `https://`
        - **Rate Limiting**: Some websites may limit requests
        - **Content Quality**: Not all websites provide extractable content
        
        **GitHub URL Examples:**
        - ✅ `https://github.com/kubeflow/trainer/blob/master/README.md`
        - ✅ `https://github.com/username/repo-name` (auto-finds README.md)
        - ✅ `https://github.com/username/repo/blob/main/docs/guide.md`
        """)

def is_valid_url_flexible(url: str) -> bool:
    """Check if URL is valid, handling both with and without protocols"""
    try:
        # If URL doesn't start with protocol, add https:// for validation
        if not url.startswith(('http://', 'https://')):
            url = f"https://{url}"
        
        # Parse the URL
        from urllib.parse import urlparse
        result = urlparse(url)
        
        # Check if it has a valid domain (netloc) and contains a dot
        return bool(result.netloc and '.' in result.netloc and len(result.netloc.split('.')[0]) > 0)
    except Exception:
        return False


def render_url_file_upload_input():
    """Render the URL file upload interface"""
    st.markdown("Upload a text or CSV file containing a list of URLs to process.")

    # File upload option for bulk URLs
    uploaded_file = st.file_uploader(
        "Upload URL list file",
        type=['txt', 'csv'],
        help="Upload a text file with one URL per line (each line will be validated as a URL), or a CSV file with URLs in the first column"
    )
    
    # Process URLs from file if uploaded
    if uploaded_file:
        try:
            if uploaded_file.name.endswith('.csv'):
                # Handle CSV file
                df = pd.read_csv(uploaded_file)
                if len(df.columns) > 0:
                    urls_from_file = df.iloc[:, 0].dropna().astype(str).tolist()
                else:
                    st.error("❌ CSV file appears to be empty")
                    return
            else:
                # Handle text file
                content = uploaded_file.read().decode('utf-8')
                # Process all non-empty lines and validate each as a URL
                all_lines = [line.strip() for line in content.split('\n') if line.strip()]
                
                # Validate URLs and filter out invalid ones
                valid_urls = []
                invalid_lines = []
                
                for line in all_lines:
                    # Validate URL using flexible validation
                    if is_valid_url_flexible(line):
                        # Add https:// if no protocol specified
                        url = line
                        if not url.startswith(('http://', 'https://')):
                            url = f"https://{url}"
                        valid_urls.append(url)
                    else:
                        invalid_lines.append(line)
                
                urls_from_file = valid_urls
                
                # Show processing summary
                if invalid_lines:
                    st.warning(f"⚠️ Found {len(invalid_lines)} invalid lines that will be skipped")
                    with st.expander("Show invalid lines"):
                        for line in invalid_lines:
                            st.code(line)
                if urls_from_file:
                    st.success(f"📄 Found {len(urls_from_file)} valid URLs to process")
            
            if urls_from_file:
                st.success(f"📄 Loaded {len(urls_from_file)} URLs from file")
                st.markdown("**URLs from file:**")
                for url in urls_from_file[:10]:  # Show first 10
                    st.code(url)
                if len(urls_from_file) > 10:
                    st.caption(f"... and {len(urls_from_file) - 10} more URLs")
                
                if st.button("🚀 Process All URLs", type="primary"):
                    if urls_from_file:
                        process_bulk_web_urls(urls_from_file)
                    else:
                        st.warning("⚠️ No URLs to process")
            else:
                st.error("❌ No valid URLs found in the uploaded file")
        
        except Exception as e:
            st.error(f"❌ Error reading file: {str(e)}")

    # Show processing tips
    with st.expander("💡 File Format Tips"):
        st.markdown("""
        **Supported File Formats:**
        - **Text files (.txt)**: One URL per line
          - Each line is validated as a URL
          - Invalid lines are automatically skipped
          - URLs without protocol (http/https) will have https:// added automatically
        - **CSV files (.csv)**: URLs in the first column
        
        **Processing Limits:**
        - Maximum 50 URLs per batch (to avoid overwhelming the system)
        - Each URL is processed sequentially for stability
        - Duplicate URLs are automatically skipped
        - Processing time depends on website response times
        """)


def create_chunks_from_content(content: str) -> List[str]:
    """Create chunks from content using existing chunking logic"""
    # Use existing chunk size settings
    CHUNK_SIZE = 3000
    CHUNK_OVERLAP = 600
    
    if not content or not content.strip():
        return []
    
    chunks = []
    text_length = len(content)
    
    if text_length <= CHUNK_SIZE:
        chunks.append(content.strip())
    else:
        start = 0
        while start < text_length:
            end = start + CHUNK_SIZE
            
            # Try to break at word boundary
            if end < text_length:
                # Look for the last space within the chunk
                last_space = content.rfind(' ', start, end)
                if last_space > start:
                    end = last_space
            
            chunk = content[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position with overlap
            start = end - CHUNK_OVERLAP
            if start >= text_length:
                break
    
    return chunks


def create_embeddings_from_chunks(chunks: List[str]) -> Any:
    """Create embeddings from chunks using existing embedding system"""
    try:
        # Use existing embedding model loading logic
        model = get_local_embedding_model()
        if model is None:
            return None
        
        # Create embeddings in batches
        batch_size = 32
        all_embeddings = []
        
        for i in range(0, len(chunks), batch_size):
            batch = chunks[i:i + batch_size]
            try:
                batch_embeddings = model.encode(batch, convert_to_numpy=True, show_progress_bar=False)
                all_embeddings.extend(batch_embeddings)
            except Exception as e:
                st.warning(f"Error encoding batch {i//batch_size + 1}: {str(e)}")
                # Create dummy embeddings for this batch
                import numpy as np
                dummy_embeddings = np.random.rand(len(batch), 384).astype(np.float32)
                all_embeddings.extend(dummy_embeddings)
        
        return all_embeddings
        
    except Exception as e:
        st.warning(f"Error creating embeddings: {str(e)}")
        # Return dummy embeddings as fallback
        import numpy as np
        return [np.random.rand(384).astype(np.float32) for _ in chunks]


def validate_uploaded_files(uploaded_files: List[Any]) -> List[Any]:
    """Validate uploaded files and return valid ones"""
    valid_files = []
    
    for file in uploaded_files:
        if file.size > 50 * 1024 * 1024:  # 50MB limit
            st.error(f"❌ File '{file.name}' is too large ({format_file_size(file.size)}). Maximum size is 50MB.")
        else:
            valid_files.append(file)
    
    return valid_files


def process_uploaded_files(files: List[Any]) -> int:
    """Process uploaded files with individual state tracking and interruption handling"""
    if not files:
        return 0
    
    from .utils import mark_upload_success, mark_upload_failed
    
    successful_files = 0
    
    for file in files:
        file_size_mb = file.size / (1024 * 1024) if hasattr(file, 'size') else 0
        file_id = f"{file.name}_{file.size}"
        
        st.markdown(f"### 📄 Processing: `{file.name}` ({file_size_mb:.1f}MB)")
        
        # Show performance estimate for large files
        if file_size_mb > 5:
            estimated_time = file_size_mb * 15  # Rough estimate: 15 seconds per MB
            st.info(f"⏱️ Large file detected. Estimated processing time: ~{estimated_time:.0f} seconds")
    
        # Add interruption warning
        st.warning("⚠️ **Don't switch models during upload** - this will interrupt processing")
        
        # Create progress container
        progress_container = st.container()
        
        with progress_container:
            # Initialize progress tracking with timing
            progress_bar = st.progress(0)
            status_text = st.empty()
            timing_text = st.empty()
            metrics_col1, metrics_col2, metrics_col3 = st.columns(3)
            
            try:
                start_time = time.time()
                
                # Step 1: Read file content
                step_start = time.time()
                progress_bar.progress(0.1)
                status_text.text("📥 Reading file content...")
                
                # Check for interruption during processing
                if file_id not in st.session_state.currently_uploading:
                    st.session_state.upload_interrupted = True
                    raise InterruptedError(f"Upload interrupted for {file.name}")
                
                content = read_file_content(file)
                if not content:
                    st.error(f"❌ Could not extract content from {file.name}")
                    mark_upload_failed(file_id)
                    continue
                
                step_time = time.time() - step_start
                timing_text.text(f"⏱️ File reading: {step_time:.1f}s")
                
                actual_size_mb = len(content.encode('utf-8')) / (1024 * 1024)
                
                with metrics_col1:
                    st.metric("📊 Content Size", f"{actual_size_mb:.2f} MB")
                
                # Step 2: Create chunks with optimization  
                step_start = time.time()
                progress_bar.progress(0.3)
                status_text.text("✂️ Creating optimized chunks...")
                
                # Check for interruption
                if file_id not in st.session_state.currently_uploading:
                    st.session_state.upload_interrupted = True
                    raise InterruptedError(f"Upload interrupted for {file.name}")
                
                chunks = create_enhanced_chunks(content)
                
                # Apply large file optimizations
                from .config import LARGE_FILE_THRESHOLD_MB, MAX_CHUNKS_PER_FILE
                if actual_size_mb > LARGE_FILE_THRESHOLD_MB and len(chunks) > MAX_CHUNKS_PER_FILE:
                    status_text.text(f"🚀 Optimizing {len(chunks)} chunks for performance...")
                    chunks = optimize_chunks_for_large_files(chunks)
                
                step_time = time.time() - step_start
                timing_text.text(f"⏱️ Chunking: {step_time:.1f}s")
                
                with metrics_col2:
                    st.metric("📦 Chunks Created", f"{len(chunks)}")
                
                # Step 3: Generate embeddings using optimized batch processing
                step_start = time.time()
                progress_bar.progress(0.5)
                status_text.text("🧠 Generating embeddings (optimized)...")
                
                # Check for interruption before expensive operation
                if file_id not in st.session_state.currently_uploading:
                    st.session_state.upload_interrupted = True
                    raise InterruptedError(f"Upload interrupted for {file.name}")
                
                embeddings, embedding_errors = generate_embeddings_batch_optimized(
                    chunks, 
                    progress_bar, 
                    status_text,
                    start_progress=0.5,
                    file_id=file_id  # Pass file_id for interruption checking
                )
                
                step_time = time.time() - step_start
                timing_text.text(f"⏱️ Embeddings: {step_time:.1f}s")
                
                with metrics_col3:
                    embedding_quality = max(0, (len(chunks) - embedding_errors) / len(chunks) * 100)
                    st.metric("🎯 Quality", f"{embedding_quality:.0f}%")
                
                # Step 4: Store document and embeddings in VectorIO database
                step_start = time.time()
                progress_bar.progress(0.8)
                status_text.text("💾 Storing in VectorIO database...")
                
                # Prepare metadata for VectorIO
                vector_metadata = []
                for i, chunk in enumerate(chunks):
                    vector_metadata.append({
                        "document_id": f"file_{file.name}_{i}",
                        "document_name": file.name,
                        "chunk_index": i,
                        "chunk_content": chunk,
                        "file_type": "FILE",
                        "file_size_mb": actual_size_mb,
                        "upload_time": time.strftime("%Y-%m-%d %H:%M:%S"),
                        "chunk_count": len(chunks),
                        "character_count": len(content)
                    })
                
                # Store embeddings in LlamaStack VectorIO database
                vectorio_success = False
                try:
                    if 'llamastack_client' in st.session_state:
                        vectorio_success = st.session_state.llamastack_client.store_embeddings_in_vector_db(
                            embeddings=embeddings,
                            metadata=vector_metadata,
                            vector_db_id="faiss"  # Use the FAISS database configured in llamastack-config.yaml
                        )
                        if vectorio_success:
                            print(f"✅ Successfully stored embeddings in VectorIO for {file.name}")
                        else:
                            print(f"⚠️ Failed to store in VectorIO, falling back to session state for {file.name}")
                except Exception as e:
                    print(f"⚠️ VectorIO storage failed for {file.name}: {e}")
                
                # Store document data with backup for session state persistence (fallback)
                doc_data = {
                    'name': file.name,
                    'content': content,
                    'chunks': chunks,
                    'embeddings': embeddings,
                    'processing_time': time.time() - start_time,
                    'chunk_count': len(chunks),
                    'embedding_errors': embedding_errors,
                    'file_size_mb': actual_size_mb,
                    'vectorio_stored': vectorio_success
                }
                
                # Check for duplicates before adding
                if 'uploaded_documents' not in st.session_state:
                    st.session_state.uploaded_documents = []
                
                # Check if document with same name and size already exists
                file_id = f"{file.name}_{file.size}"
                existing_docs = [doc for doc in st.session_state.uploaded_documents 
                               if doc.get('name') == file.name and 
                               doc.get('file_size_mb', 0) == actual_size_mb]
                
                if existing_docs:
                    st.warning(f"⚠️ File '{file.name}' already exists in your documents")
                    mark_upload_success(file_id)  # Mark as processed to avoid retry
                    continue
                
                # Add to documents array (only one array to avoid duplication)
                st.session_state.uploaded_documents.append(doc_data)
                
                # Final progress
                progress_bar.progress(1.0)
                total_time = time.time() - start_time
                
                status_text.text("✅ Processing complete!")
                timing_text.text(f"⏱️ Total time: {total_time:.1f}s ({actual_size_mb/total_time:.1f} MB/s)")
                
                # Show performance summary
                storage_method = "VectorIO + Session" if vectorio_success else "Session State (fallback)"
                st.success(f"""
                📄 **{file.name}** processed successfully!
                - **Size:** {actual_size_mb:.1f}MB → {len(chunks)} chunks
                - **Quality:** {embedding_quality:.0f}% embeddings successful
                - **Storage:** {storage_method}
                - **Speed:** {total_time:.1f}s ({actual_size_mb/total_time:.1f} MB/s)
                """)
                
                # Mark this file as successfully processed
                mark_upload_success(file_id)
                successful_files += 1
                
            except InterruptedError as e:
                st.error(f"⚠️ Upload interrupted: {e}")
                mark_upload_failed(file_id)
                st.info("💡 File marked for retry - try uploading again after model operations complete")
                
            except Exception as e:
                st.error(f"❌ Error processing {file.name}: {str(e)}")
                mark_upload_failed(file_id)
                st.info("💡 Check file format and try again")
    
    return successful_files


def render_document_library() -> None:
    """Render the document library interface with original layout"""
    # Try to restore data first in case it was lost
    restore_documents_data()
    
    # Get total document count
    doc_count = get_document_count()
    
    # Debug: Check document state at render time
    if doc_count == 0:
        print(f"⚠️ Document library render: no documents found")
    else:
        print(f"✅ Document library render: found {doc_count} documents")
    
    st.markdown("---")
    
    # Ensure session state exists
    if 'uploaded_documents' not in st.session_state:
        st.session_state.uploaded_documents = []
    
    # Use only uploaded_documents to avoid duplication
    all_documents = st.session_state.uploaded_documents
    
    # Count documents for the header
    total_size_mb = sum(doc.get("file_size_mb", 0) for doc in all_documents) if all_documents else 0
    
    # Create collapsible expander with document count
    header_text = f"📁 Your Documents ({doc_count})"
    if doc_count > 0:
        header_text += f" • {total_size_mb:.1f}MB"
    
    with st.expander(header_text, expanded=doc_count > 0):  # Auto-expand when documents exist
        if all_documents and doc_count > 0:
            # Show document summary
            total_chunks = sum(doc.get("chunk_count", 0) for doc in all_documents)
            total_chars = sum(len(doc.get("content", "")) for doc in all_documents)
            
            # Quick stats
            st.info(f"📊 {doc_count} documents • {total_size_mb:.1f}MB total • {total_chunks} chunks")
            
            # Search and filter options
            col1, col2 = st.columns([2, 1])
            with col1:
                search_term = st.text_input("🔍 Search documents", placeholder="Search by name, content, or URL...", key="doc_search")
            with col2:
                filter_type = st.selectbox("📁 Filter by type", ["All", "WEB", "PDF", "TXT", "DOCX"], key="doc_filter")
            
            # Filter documents based on search and filter
            filtered_documents = all_documents
            if search_term:
                search_lower = search_term.lower()
                filtered_documents = [
                    doc for doc in filtered_documents
                    if (search_lower in doc.get('name', '').lower() or
                        search_lower in doc.get('content', '').lower() or
                        search_lower in doc.get('source_url', '').lower())
                ]
            
            if filter_type != "All":
                filtered_documents = [
                    doc for doc in filtered_documents
                    if doc.get('file_type') == filter_type
                ]
            
            # Show filtered count
            if len(filtered_documents) != len(all_documents):
                st.info(f"🔍 Showing {len(filtered_documents)} of {len(all_documents)} documents")
            
            # Embedding Quality Summary
            real_embeddings = sum(doc.get('chunk_count', 0) - doc.get('embedding_errors', 0) for doc in all_documents)
            total_embeddings = sum(doc.get('chunk_count', 0) for doc in all_documents)
            embedding_quality = (real_embeddings / total_embeddings * 100) if total_embeddings > 0 else 0
            
            if embedding_quality > 80:
                quality_status = "Excellent"
                st.success(f"🎯 Embedding Quality: {embedding_quality:.1f}% ({real_embeddings}/{total_embeddings}) - {quality_status}")
            elif embedding_quality > 50:
                quality_status = "Good"
                st.warning(f"🎯 Embedding Quality: {embedding_quality:.1f}% ({real_embeddings}/{total_embeddings}) - {quality_status}")
            else:
                quality_status = "Needs Improvement"
                st.error(f"🎯 Embedding Quality: {embedding_quality:.1f}% ({real_embeddings}/{total_embeddings}) - {quality_status}")
            
            # Clear all button
            col1, col2 = st.columns([1, 4])
            with col1:
                if st.button("🗑️ Clear All", type="secondary", help="Remove all documents"):
                    clear_all_documents()
                    st.rerun()
            
            # Detailed Performance Results Table (in dropdown)
            with st.expander("📊 Detailed Performance Results", expanded=False):
                # Calculate additional metrics
                avg_chunk_size = total_chars // total_chunks if total_chunks > 0 else 0
                avg_file_size = total_size_mb / doc_count if doc_count > 0 else 0
                total_processing_time = sum(doc.get('processing_time', 0) for doc in all_documents)
                avg_processing_speed = total_size_mb / total_processing_time if total_processing_time > 0 else 0
                
                # Create comprehensive performance table
                performance_data = {
                    "Metric": [
                        "📄 Documents Processed",
                        "📦 Total File Size",
                        "📝 Characters Extracted", 
                        "✂️ Text Chunks Created",
                        "🧠 Embeddings Generated",
                        "🎯 Embedding Quality",
                        "⏱️ Total Processing Time",
                        "🚀 Average Processing Speed",
                        "📏 Average Chunk Size",
                        "📊 Average File Size",
                        "🔢 Chunks per Document",
                        "💯 Success Rate",
                        "🔗 Characters per Document",
                        "✨ Status"
                    ],
                    "Value": [
                        f"{doc_count:,}",
                        f"{total_size_mb:.2f} MB",
                        f"{total_chars:,}",
                        f"{total_chunks:,}",
                        f"{real_embeddings:,} / {total_embeddings:,}",
                        f"{embedding_quality:.1f}% ({quality_status})",
                        f"{total_processing_time:.1f} seconds",
                        f"{avg_processing_speed:.2f} MB/s" if avg_processing_speed > 0 else "N/A",
                        f"{avg_chunk_size:,} characters",
                        f"{avg_file_size:.2f} MB",
                        f"{total_chunks // doc_count if doc_count > 0 else 0:,}",
                        f"{(doc_count / len(all_documents) * 100):.1f}%" if all_documents else "N/A",
                        f"{total_chars // doc_count:,}" if doc_count > 0 else "N/A",
                        "🟢 Ready for Q&A and Search"
                    ],
                    "Category": [
                        "📋 Overview",
                        "📋 Overview", 
                        "📋 Overview",
                        "📋 Overview",
                        "🧠 AI Processing",
                        "🧠 AI Processing",
                        "⚡ Performance",
                        "⚡ Performance",
                        "📊 Statistics",
                        "📊 Statistics",
                        "📊 Statistics",
                        "📊 Statistics",
                        "📊 Statistics",
                        "✅ System"
                    ]
                }
                
                # Create and display the dataframe
                performance_df = pd.DataFrame(performance_data)
                
                # Style the table
                st.markdown("**🔍 Comprehensive Performance Analysis**")
                st.dataframe(
                    performance_df,
                    use_container_width=True,
                    hide_index=True,
                    height=500
                )
                
                # Individual document breakdown
                if len(all_documents) > 1:
                    st.markdown("---")
                    st.markdown("**📄 Individual Document Performance**")
                    
                    doc_performance = []
                    for doc in all_documents:
                        chunk_count = doc.get('chunk_count', 0)
                        embedding_errors = doc.get('embedding_errors', 0)
                        doc_quality = ((chunk_count - embedding_errors) / chunk_count * 100) if chunk_count > 0 else 0
                        
                        doc_performance.append({
                            "Document": doc['name'],
                            "Size (MB)": f"{doc.get('file_size_mb', 0):.2f}",
                            "Characters": f"{len(doc.get('content', '')):,}",
                            "Chunks": f"{chunk_count:,}",
                            "Processing Time (s)": f"{doc.get('processing_time', 0):.1f}",
                            "Embedding Quality": f"{doc_quality:.1f}%",
                            "Speed (MB/s)": f"{doc.get('file_size_mb', 0) / max(doc.get('processing_time', 1), 0.1):.2f}"
                        })
                    
                    doc_df = pd.DataFrame(doc_performance)
                    st.dataframe(
                        doc_df,
                        use_container_width=True,
                        hide_index=True
                    )
            
            # Simple document list
            for i, doc in enumerate(filtered_documents):
                with st.expander(f"📄 {doc['name']}", expanded=False):
                    # Basic document info
                    col1, col2 = st.columns([1, 1])
                    with col1:
                        st.write(f"**Size:** {doc.get('file_size_mb', 0):.2f} MB")
                        st.write(f"**Chunks:** {doc.get('chunk_count', 0)}")
                        st.write(f"**Characters:** {len(doc.get('content', '')):,}")
                    
                    with col2:
                        # Show MCP server information
                        if doc.get('extraction_method'):
                            extraction_method = doc.get('extraction_method')
                            method_map = {
                                'mcp_just-every': '📦 Just-Every MCP'
                            }
                            method_display = method_map.get(extraction_method, extraction_method)
                            st.info(f"🔧 **MCP Server Used:** {method_display}")
                            
                            # Show source URL
                            if doc.get('source_url'):
                                st.write(f"**🔗 Source:** {doc.get('source_url')}")
                        
                        # Show content preview
                        content = doc.get('content', '')
                        if not content and 'llamastack_client' in st.session_state:
                            # Try to retrieve content from VectorIO if not in session state
                            try:
                                # Search for chunks from this document in VectorIO
                                doc_name = doc.get('name', '')
                                if doc_name:
                                    # Use a broad search to get chunks from this document
                                    vectorio_results = st.session_state.llamastack_client.search_similar_vectors(
                                        query_text=doc_name,
                                        vector_db_id="faiss",
                                        top_k=50
                                    )
                                    
                                    # Handle VectorIO response format (dict with 'chunks' key)
                                    if isinstance(vectorio_results, dict) and 'chunks' in vectorio_results:
                                        vectorio_results = vectorio_results['chunks']
                                    
                                    # Filter results to only include chunks from this document
                                    doc_chunks = []
                                    for result in vectorio_results:
                                        metadata = result.get('metadata', {})
                                        result_doc_name = metadata.get('document_name', '')
                                        
                                        # Try multiple matching strategies
                                        if (result_doc_name == doc_name or 
                                            result_doc_name.endswith(doc_name) or
                                            doc_name in result_doc_name or
                                            result_doc_name.replace('Web Content from ', '').replace('https://', '').replace('http://', '').replace('/', '') == doc_name.replace('🌐 ', '').replace('📄 ', '')):
                                            chunk_content = result.get('content', '')
                                            if chunk_content:
                                                doc_chunks.append(chunk_content)
                                    
                                    # Combine chunks to reconstruct content
                                    if doc_chunks:
                                        content = ' '.join(doc_chunks)
                                        # Update the document with retrieved content
                                        doc['content'] = content
                            except Exception as e:
                                print(f"⚠️ Could not retrieve content from VectorIO for {doc.get('name', 'Unknown')}: {e}")
                        
                        if content:
                            preview = content[:200].replace('\n', ' ').strip()
                            if len(content) > 200:
                                preview += "..."
                            st.caption(f"**Preview:** {preview}")
                        else:
                            st.caption("**Preview:** Content not available in session state or VectorIO")
                    
                    # Action buttons
                    col1, col2, col3 = st.columns([1, 1, 1])
                    with col1:
                        if st.button(f"👁️ View Content", key=f"view_{i}"):
                            st.session_state[f"show_content_{i}"] = not st.session_state.get(f"show_content_{i}", False)
                            st.rerun()
                    
                    with col2:
                        if st.button(f"📊 View Chunks", key=f"chunks_{i}"):
                            st.session_state[f"show_chunks_{i}"] = not st.session_state.get(f"show_chunks_{i}", False)
                            st.rerun()
                    
                    with col3:
                        if st.button(f"🗑️ Remove", key=f"remove_{i}"):
                            remove_document(i)
                            st.rerun()
                    
                    # Show content if requested
                    if st.session_state.get(f"show_content_{i}", False):
                        st.markdown("---")
                        st.markdown("### 📄 Raw Extracted Content")
                        
                        # Show extraction method and statistics
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("Method", doc.get('extraction_method', 'unknown'))
                        with col2:
                            st.metric("Characters", f"{len(doc.get('content', '')):,}")
                        with col3:
                            st.metric("Chunks", doc.get('chunk_count', 0))
                        
                        # Show additional metadata for web documents
                        if doc.get('file_type') == 'WEB':
                            if doc.get('source_url'):
                                st.info(f"**Source:** {doc.get('source_url')}")
                            if doc.get('metadata'):
                                with st.expander("📊 Extraction Metadata", expanded=False):
                                    st.json(doc.get('metadata', {}))
                        
                        # Show content in a scrollable area
                        content = doc.get('content', '')
                        
                        # If content is not available, try to retrieve from VectorIO
                        if not content and 'llamastack_client' in st.session_state:
                            try:
                                doc_name = doc.get('name', '')
                                if doc_name:
                                    # Get all chunks for this document from VectorIO
                                    vectorio_results = st.session_state.llamastack_client.search_similar_vectors(
                                        query_text=doc_name,
                                        vector_db_id="faiss",
                                        top_k=100
                                    )
                                    
                                    # Handle VectorIO response format (dict with 'chunks' key)
                                    if isinstance(vectorio_results, dict) and 'chunks' in vectorio_results:
                                        vectorio_results = vectorio_results['chunks']
                                    
                                    # Filter and combine chunks
                                    doc_chunks = []
                                    for result in vectorio_results:
                                        metadata = result.get('metadata', {})
                                        result_doc_name = metadata.get('document_name', '')
                                        
                                        # Try multiple matching strategies
                                        if (result_doc_name == doc_name or 
                                            result_doc_name.endswith(doc_name) or
                                            doc_name in result_doc_name or
                                            result_doc_name.replace('Web Content from ', '').replace('https://', '').replace('http://', '').replace('/', '') == doc_name.replace('🌐 ', '').replace('📄 ', '')):
                                            chunk_content = result.get('content', '')
                                            if chunk_content:
                                                doc_chunks.append(chunk_content)
                                    
                                    if doc_chunks:
                                        content = ' '.join(doc_chunks)
                                        doc['content'] = content
                                        st.success(f"✅ Retrieved {len(doc_chunks)} chunks from VectorIO database")
                            except Exception as e:
                                st.error(f"❌ Could not retrieve content from VectorIO: {e}")
                        
                        if content:
                            if len(content) > 5000:
                                st.markdown("**Content Preview (first 5000 characters):**")
                                st.text_area("Raw Content", content[:5000], height=300, key=f"content_preview_{i}")
                                st.info(f"Content truncated. Full content has {len(content):,} characters.")
                            else:
                                st.text_area("Raw Content", content, height=300, key=f"content_full_{i}")
                        else:
                            st.warning("⚠️ Content not available in session state or VectorIO database")
                    
                    # Show chunks if requested
                    if st.session_state.get(f"show_chunks_{i}", False):
                        st.markdown("---")
                        st.markdown("### 🧩 Document Chunks")
                        chunks = doc.get('chunks', [])
                        
                        # If chunks are not available, try to retrieve from VectorIO
                        if not chunks and 'llamastack_client' in st.session_state:
                            try:
                                doc_name = doc.get('name', '')
                                if doc_name:
                                    # Get all chunks for this document from VectorIO
                                    vectorio_results = st.session_state.llamastack_client.search_similar_vectors(
                                        query_text=doc_name,
                                        vector_db_id="faiss",
                                        top_k=100
                                    )
                                    
                                    # Handle VectorIO response format (dict with 'chunks' key)
                                    if isinstance(vectorio_results, dict) and 'chunks' in vectorio_results:
                                        vectorio_results = vectorio_results['chunks']
                                    
                                    # Filter and extract chunks
                                    doc_chunks = []
                                    for result in vectorio_results:
                                        metadata = result.get('metadata', {})
                                        result_doc_name = metadata.get('document_name', '')
                                        
                                        # Try multiple matching strategies
                                        if (result_doc_name == doc_name or 
                                            result_doc_name.endswith(doc_name) or
                                            doc_name in result_doc_name or
                                            result_doc_name.replace('Web Content from ', '').replace('https://', '').replace('http://', '').replace('/', '') == doc_name.replace('🌐 ', '').replace('📄 ', '')):
                                            chunk_content = result.get('content', '')
                                            if chunk_content:
                                                doc_chunks.append(chunk_content)
                                    
                                    if doc_chunks:
                                        chunks = doc_chunks
                                        doc['chunks'] = chunks
                                        st.success(f"✅ Retrieved {len(chunks)} chunks from VectorIO database")
                            except Exception as e:
                                st.error(f"❌ Could not retrieve chunks from VectorIO: {e}")
                        
                        if chunks:
                            for j, chunk in enumerate(chunks):
                                with st.expander(f"Chunk {j+1} ({len(chunk)} chars)", expanded=False):
                                    st.text(chunk)
                        else:
                            st.info("No chunks available for this document in session state or VectorIO database.")
        
        else:
            st.info("🔍 No documents uploaded yet. Upload some documents above to get started!")


def clear_all_documents() -> None:
    """Clear all uploaded documents"""
    st.session_state.uploaded_documents = []
    st.session_state.documents = []  # Also clear old documents for chat
    st.session_state.chat_history = []  # Clear chat history too

    # Clear backup as well
    if '_backup_uploaded_documents' in st.session_state:
        st.session_state['_backup_uploaded_documents'] = []
    if '_backup_documents' in st.session_state:
        st.session_state['_backup_documents'] = []
    
    print("🗑️ Cleared all documents and backups")


def remove_document(idx: int) -> None:
    """Remove a specific document by index from the combined document list"""
    # Get the combined document list
    all_documents = st.session_state.documents + st.session_state.uploaded_documents
    
    if idx >= len(all_documents):
        print(f"❌ Invalid document index: {idx}")
        return
    
    # Get the document to remove
    doc_to_remove = all_documents[idx]
    doc_name = doc_to_remove['name']
    
    # Remove from appropriate collection
    if doc_to_remove in st.session_state.documents:
        st.session_state.documents.remove(doc_to_remove)
        print(f"🗑️ Removed document from 'documents': {doc_name}")
    elif doc_to_remove in st.session_state.uploaded_documents:
        st.session_state.uploaded_documents.remove(doc_to_remove)
        print(f"🗑️ Removed document from 'uploaded_documents': {doc_name}")
    
    # Update backup
    backup_documents_data()
    
    print(f"🗑️ Removed document: {doc_name}")


def has_documents() -> bool:
    """Check if any documents are uploaded"""
    if 'uploaded_documents' in st.session_state:
        return len(st.session_state.uploaded_documents) > 0
    return False


def get_document_count() -> int:
    """Get the total number of uploaded documents"""
    if 'uploaded_documents' in st.session_state:
        return len(st.session_state.uploaded_documents) 
    return 0 